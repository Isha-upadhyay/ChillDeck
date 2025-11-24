# app/services/export_service.py

from pathlib import Path
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
import uuid
import json
from app.core.logger import logger


class ExportService:
    """
    Handles exporting slides into different formats:
    - PPTX (PowerPoint)
    - PDF
    - Markdown
    """

    def __init__(self, export_dir: str = "exports"):
        self.export_dir = Path(export_dir)
        self.export_dir.mkdir(parents=True, exist_ok=True)

    def _generate_file_path(self, suffix: str, topic: str = "slides") -> Path:
        """Generate a file path for export"""
        # Clean topic name for filename
        safe_topic = "".join(c for c in topic if c.isalnum() or c in (' ', '-', '_')).rstrip()[:50]
        safe_topic = safe_topic.replace(' ', '_') if safe_topic else "slides"
        filename = f"{safe_topic}_{uuid.uuid4().hex[:8]}.{suffix}"
        return self.export_dir / filename

    # --------------------------
    # EXPORT: PPTX
    # --------------------------
    def export_pptx(self, slides: List[Dict[str, Any]], topic: str = "Presentation") -> str:
        """Export slides to PowerPoint format."""
        try:
            pptx_path = self._generate_file_path("pptx", topic)
            prs = Presentation()
            
            # Set slide dimensions (16:9 aspect ratio)
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)

            for idx, slide_data in enumerate(slides, start=1):
                # Use blank layout
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Get slide content and theme
                title = slide_data.get("title") or slide_data.get("heading", f"Slide {idx}")
                bullets = slide_data.get("bullets") or slide_data.get("points", [])
                theme = slide_data.get("design", {}).get("theme", "corporate")
                
                # Apply theme colors (basic implementation)
                theme_colors = self._get_theme_colors(theme)
                
                # Add title
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                height = Inches(1)
                title_box = slide.shapes.add_textbox(left, top, width, height)
                title_frame = title_box.text_frame
                title_frame.text = title
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(44)
                title_para.font.bold = True
                title_para.alignment = PP_ALIGN.LEFT
                if theme_colors.get("title_color"):
                    title_para.font.color.rgb = theme_colors["title_color"]
                
                # Add bullets
                bullet_top = Inches(2)
                bullet_left = Inches(1)
                bullet_width = Inches(8.5)
                
                for bullet_idx, bullet_text in enumerate(bullets[:6]):  # Limit to 6 bullets
                    bullet_height = Inches(0.6)
                    bullet_box = slide.shapes.add_textbox(
                        bullet_left, 
                        bullet_top + (bullet_idx * Inches(0.7)), 
                        bullet_width, 
                        bullet_height
                    )
                    bullet_frame = bullet_box.text_frame
                    bullet_frame.text = f"• {bullet_text}"
                    bullet_para = bullet_frame.paragraphs[0]
                    bullet_para.font.size = Pt(24)
                    bullet_para.alignment = PP_ALIGN.LEFT
                    if theme_colors.get("text_color"):
                        bullet_para.font.color.rgb = theme_colors["text_color"]

            prs.save(str(pptx_path))
            logger.info(f"Exported PPTX to: {pptx_path}")
            return str(pptx_path)
        except Exception as e:
            logger.error(f"Error exporting PPTX: {e}")
            raise

    def _get_theme_colors(self, theme: str) -> Dict[str, Any]:
        """Get color scheme for a theme"""
        try:
            from pptx.dml.color import RGBColor
        except ImportError:
            # Fallback if RGBColor not available
            RGBColor = None
        
        themes = {
            "corporate": {
                "title_color": RGBColor(30, 64, 175) if RGBColor else None,  # Blue
                "text_color": RGBColor(30, 30, 30) if RGBColor else None,    # Dark gray
            },
            "dark": {
                "title_color": RGBColor(255, 255, 255) if RGBColor else None,  # White
                "text_color": RGBColor(200, 200, 200) if RGBColor else None,   # Light gray
            },
            "modern": {
                "title_color": RGBColor(147, 51, 234) if RGBColor else None,  # Purple
                "text_color": RGBColor(50, 50, 50) if RGBColor else None,     # Dark gray
            },
            "tech": {
                "title_color": RGBColor(22, 163, 74) if RGBColor else None,  # Green
                "text_color": RGBColor(30, 30, 30) if RGBColor else None,    # Dark gray
            },
            "cute": {
                "title_color": RGBColor(236, 72, 153) if RGBColor else None,  # Pink
                "text_color": RGBColor(50, 50, 50) if RGBColor else None,     # Dark gray
            },
            "minimal": {
                "title_color": RGBColor(0, 0, 0) if RGBColor else None,       # Black
                "text_color": RGBColor(50, 50, 50) if RGBColor else None,    # Dark gray
            },
        }
        return themes.get(theme, themes["corporate"])

    # --------------------------
    # EXPORT: PDF
    # --------------------------
    def export_pdf(self, slides: List[Dict[str, Any]], topic: str = "Presentation") -> str:
        """Export slides to PDF format."""
        try:
            pdf_path = self._generate_file_path("pdf", topic)
            doc = SimpleDocTemplate(str(pdf_path), pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            # Custom styles
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor='#1e40af',
                spaceAfter=30,
                alignment=1  # Center
            )
            
            bullet_style = ParagraphStyle(
                'CustomBullet',
                parent=styles['Normal'],
                fontSize=12,
                leftIndent=20,
                spaceAfter=10
            )

            for idx, slide_data in enumerate(slides, start=1):
                if idx > 1:
                    story.append(PageBreak())
                
                # Add title
                title = slide_data.get("title") or slide_data.get("heading", f"Slide {idx}")
                story.append(Paragraph(title, title_style))
                story.append(Spacer(1, 0.3*inch))
                
                # Add bullets
                bullets = slide_data.get("bullets") or slide_data.get("points", [])
                for bullet in bullets:
                    story.append(Paragraph(f"• {bullet}", bullet_style))
                    story.append(Spacer(1, 0.1*inch))
                
                # Add notes if available
                notes = slide_data.get("notes", "")
                if notes:
                    story.append(Spacer(1, 0.2*inch))
                    story.append(Paragraph(f"<i>Notes: {notes}</i>", styles['Italic']))

            doc.build(story)
            logger.info(f"Exported PDF to: {pdf_path}")
            return str(pdf_path)
        except Exception as e:
            logger.error(f"Error exporting PDF: {e}")
            raise

    # --------------------------
    # EXPORT: MARKDOWN
    # --------------------------
    def export_markdown(self, slides: List[Dict[str, Any]], topic: str = "Presentation") -> str:
        """Export slides to Markdown format."""
        try:
            md_path = self._generate_file_path("md", topic)
            
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(f"# {topic}\n\n")
                f.write(f"*Generated presentation with {len(slides)} slides*\n\n")
                f.write("---\n\n")
                
                for idx, slide_data in enumerate(slides, start=1):
                    title = slide_data.get("title") or slide_data.get("heading", f"Slide {idx}")
                    bullets = slide_data.get("bullets") or slide_data.get("points", [])
                    notes = slide_data.get("notes", "")
                    
                    f.write(f"## Slide {idx}: {title}\n\n")
                    
                    for bullet in bullets:
                        f.write(f"- {bullet}\n")
                    
                    if notes:
                        f.write(f"\n*Notes: {notes}*\n")
                    
                    f.write("\n---\n\n")
            
            logger.info(f"Exported Markdown to: {md_path}")
            return str(md_path)
        except Exception as e:
            logger.error(f"Error exporting Markdown: {e}")
            raise

    # --------------------------
    # EXPORT: JSON
    # --------------------------
    def export_json(self, slides: List[Dict[str, Any]], topic: str = "Presentation") -> str:
        """Export slides to JSON format."""
        try:
            json_path = self._generate_file_path("json", topic)
            
            export_data = {
                "topic": topic,
                "total_slides": len(slides),
                "slides": slides
            }
            
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(export_data, f, indent=2, ensure_ascii=False)
            
            logger.info(f"Exported JSON to: {json_path}")
            return str(json_path)
        except Exception as e:
            logger.error(f"Error exporting JSON: {e}")
            raise
